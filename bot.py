import os
import io
import sys
import json
import logging
import contextlib
import httpx
from io import BytesIO
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from telegram import Update
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    filters,
    ContextTypes,
)

load_dotenv()

TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")

NVIDIA_URL = "https://integrate.api.nvidia.com/v1/chat/completions"
NVIDIA_KEY = "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"

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

http_client = httpx.AsyncClient(timeout=30.0, headers={
    "Authorization": f"Bearer {NVIDIA_KEY}",
    "Content-Type": "application/json",
})

img_client = httpx.AsyncClient(timeout=60.0, follow_redirects=True)


async def ai_chat(messages: list, max_tokens: int = 1024) -> str:
    payload = {
        "model": "meta/llama-3.1-8b-instruct",
        "messages": messages,
        "max_tokens": max_tokens,
        "temperature": 0.7,
        "top_p": 0.9,
    }
    response = await http_client.post(NVIDIA_URL, json=payload)
    response.raise_for_status()
    return response.json()["choices"][0]["message"]["content"]


def create_presentation(title: str, slides_data: list) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(15, 23, 42)
    accent_color = RGBColor(59, 130, 246)
    text_color = RGBColor(255, 255, 255)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        if "title" in slide_info:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_info["title"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = accent_color

        if "content" in slide_info:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, point in enumerate(slide_info["content"]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"  {point}"
                p.font.size = Pt(24)
                p.font.color.rgb = text_color
                p.space_after = Pt(12)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def detect_presentation(text: str) -> dict:
    lower = text.lower()

    create_triggers = ["создай презентацию", "сделай презентацию", "презентация на",
                       "создать презентацию", "сделать презентацию", "нужна презентация",
                       "подготовь презентацию", "make presentation"]
    for t in create_triggers:
        if t in lower:
            return {"action": "create", "topic": text}

    edit_triggers = ["отредактируй презентацию", "измени презентацию", "редактировать",
                     "добавь слайд", "удали слайд", "измени слайд"]
    for t in edit_triggers:
        if t in lower:
            return {"action": "edit", "topic": text}

    return {"action": None}


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    await update.message.reply_text(
        "Привет! Я AI-бот с инструментами.\n\n"
        "Просто напиши мне:\n"
        "- Создай презентацию про космос\n"
        "- Найди в интернете что-нибудь\n"
        "- Напиши код для чего-нибудь\n"
        "- Переведи на английский привет\n\n"
        "/image <промпт> - сгенерировать картинку"
    )


async def handle_image(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    prompt = " ".join(context.args) if context.args else None

    if not prompt:
        await update.message.reply_text("Используй: /image <описание>")
        return

    await update.message.reply_text("Генерирую...")

    try:
        url = f"https://image.pollinations.ai/prompt/{prompt}?width=1024&height=1024&nologo=true&seed=-1"
        response = await img_client.get(url)
        response.raise_for_status()

        img_bytes = BytesIO(response.content)
        img_bytes.name = "image.png"

        await update.message.reply_photo(photo=img_bytes, caption=f"По запросу: {prompt}")
    except Exception as e:
        logger.error(f"Image error: {e}")
        await update.message.reply_text("Не удалось сгенерировать картинку.")


async def do_search(query: str) -> str:
    try:
        from duckduckgo_search import AsyncDDGS

        async with AsyncDDGS() as ddgs:
            results = []
            async for r in ddgs.text(query, max_results=5):
                results.append(f"{r['title']}\n{r['href']}\n{r['body']}")

        return "\n\n".join(results) if results else "Ничего не найдено."
    except Exception as e:
        return f"Ошибка поиска: {e}"


async def do_code(code: str) -> str:
    try:
        old_stdout = sys.stdout
        sys.stdout = buffer = io.StringIO()

        with contextlib.redirect_stdout(buffer):
            exec(code, {"__builtins__": __builtins__}, {})

        output = buffer.getvalue()
        sys.stdout = old_stdout

        return output.strip() if output.strip() else "(нет вывода)"
    except Exception as e:
        sys.stdout = old_stdout
        return f"Ошибка: {e}"


async def do_translate(text: str) -> str:
    try:
        lower = text.lower()
        target = "en"

        lang_map = {"на английский": "en", "на русский": "ru", "на испанский": "es",
                    "на французский": "fr", "на немецкий": "de", "на китайский": "zh",
                    "на японский": "ja", "на корейский": "ko", "translate to english": "en"}

        for phrase, lang in lang_map.items():
            if phrase in lower:
                target = lang
                break

        clean = text
        for phrase in ["переведи на", "перевод на", "как будет на", "translate to", "перевести на"]:
            clean = clean.replace(phrase, "").replace(phrase.upper(), "")
        clean = clean.strip()

        url = "https://api.mymemory.translated.net/get"
        params = {"q": clean, "langpair": f"ru|{target}"}

        response = await img_client.get(url, params=params)
        response.raise_for_status()
        data = response.json()

        return data["responseData"]["translatedText"]
    except Exception as e:
        return f"Ошибка перевода: {e}"


async def do_presentation(topic: str) -> BytesIO:
    prompt = f"""Создай структуру презентации на тему: {topic}
Ответ ТОЛЬКО в формате JSON без markdown:
{{"title": "Название презентации", "slides": [{{"title": "Название слайда", "content": ["Пункт 1", "Пункт 2", "Пункт 3"]}}]}}
Сделай 5-7 слайдов, каждый с 3-4 пунктами."""

    messages = [{"role": "user", "content": prompt}]
    response_text = await ai_chat(messages, max_tokens=1024)

    try:
        json_str = response_text.strip()
        if "```" in json_str:
            json_str = json_str.split("```")[1]
            if json_str.startswith("json"):
                json_str = json_str[4:]
            json_str = json_str.strip()

        data = json.loads(json_str)
        title = data.get("title", topic)
        slides = data.get("slides", [])
    except Exception as e:
        logger.error(f"JSON parse error: {e}, response: {response_text}")
        title = topic
        slides = [
            {"title": "Введение", "content": [f"Тема: {topic}", "Обзор основных понятий", "Актуальность темы"]},
            {"title": "Основная часть", "content": ["Ключевые концепции", "Примеры", "Детали"]},
            {"title": "Заключение", "content": ["Выводы", "Перспективы", "Вопросы"]}
        ]

    return create_presentation(title, slides)


async def chat_with_ai(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_message = update.message.text
    if not user_message:
        return

    lower = user_message.lower().strip()

    is_presentation = any(word in lower for word in [
        "презентация", "presentation", "слайды", "slides",
        "создай презентацию", "сделай презентацию"
    ])

    if is_presentation:
        await update.message.reply_text("Создаю презентацию...")

        topic = user_message
        for phrase in ["создай презентацию про ", "сделай презентацию про ",
                       "создай презентацию о ", "сделай презентацию о ",
                       "создай презентацию ", "сделай презентацию ",
                       "презентация на тему ", "презентация про ",
                       "презентация о ", "презентация ",
                       "создай слайды про ", "сделай слайды про ",
                       "создай слайды ", "сделай слайды ",
                       "слайды про ", "слайды "]:
            if phrase in lower:
                topic = user_message.lower().replace(phrase, "").strip()
                topic = user_message[len(phrase):].strip()
                break

        if not topic:
            topic = user_message

        try:
            pptx_bytes = await do_presentation(topic)
            pptx_bytes.name = f"{topic[:20]}.pptx"

            await update.message.reply_document(
                document=pptx_bytes,
                caption=f"Презентация: {topic}"
            )
        except Exception as e:
            logger.error(f"Presentation error: {e}", exc_info=True)
            await update.message.reply_text("Ошибка создания презентации.")
        return

    if any(word in lower for word in ["найди", "поищи", "что такое", "какой ", "какая ",
                                       "какие ", "новости", "google", "найти"]):
        await update.message.reply_text("Ищу...")
        result = await do_search(user_message)
        try:
            reply = await ai_chat([
                {"role": "system", "content": "Сформируй ответ. Отвечай на русском."},
                {"role": "user", "content": f"Вопрос: {user_message}\n\nРезультаты:\n{result}"}
            ])
            await update.message.reply_text(reply)
        except:
            await update.message.reply_text(result[:4000])
        return

    if any(word in lower for word in ["напиши код", "код для", "программа", "запусти",
                                       "выполни", "рассчитай", "посчитай", "сколько будет",
                                       "вычисли", "посчитай"]):
        await update.message.reply_text("Выполняю...")
        code = user_message
        for phrase in ["напиши код для ", "код для ", "сделай код для ",
                       "программа на python для ", "запусти код ", "выполни код ",
                       "рассчитай ", "посчитай ", "сколько будет ", "вычисли "]:
            if phrase in lower:
                code = user_message[len(phrase):].strip()
                break
        if not code:
            code = user_message
        result = await do_code(code)
        await update.message.reply_text(f"Результат:\n```\n{result}\n```", parse_mode="Markdown")
        return

    if any(word in lower for word in ["переведи", "перевод на", "как будет на", "translate"]):
        await update.message.reply_text("Перевожу...")
        result = await do_translate(user_message)
        await update.message.reply_text(f"Перевод:\n{result}")
        return

    try:
        reply = await ai_chat([
            {"role": "system", "content": "Ты полезный помощник. Отвечай кратко на русском."},
            {"role": "user", "content": user_message}
        ])
        if len(reply) > 4000:
            reply = reply[:4000] + "..."
        await update.message.reply_text(reply)
    except Exception as e:
        logger.error(f"AI error: {e}")
        await update.message.reply_text("Ошибка. Попробуй позже.")


def main() -> None:
    if not TELEGRAM_BOT_TOKEN:
        raise RuntimeError("TELEGRAM_BOT_TOKEN not set in .env")

    app = ApplicationBuilder().token(TELEGRAM_BOT_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("image", handle_image))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, chat_with_ai))

    logger.info("Bot is starting...")
    app.run_polling()


if __name__ == "__main__":
    main()
